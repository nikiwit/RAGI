"""
Generic document processing for knowledge base.
"""

import os
import re
import logging
from typing import List, Dict, Any, Optional
from langchain_core.documents import Document

logger = logging.getLogger("RAGI")

class GenericKnowledgeBaseParser:
    """Generic parser for knowledge base documents."""

    @staticmethod
    def parse_kb(text: str, source: str, filename: str) -> List[Document]:
        """
        Parse knowledge base formatted text into individual Q&A documents.

        Format expected:
        Q: Question text
        A: Answer text

        Args:
            text: Raw text content
            source: File path
            filename: File name

        Returns:
            List of Document objects
        """
        documents = []

        # Split by Q: markers to find questions
        qa_pattern = r'Q:\s*(.*?)\s*A:\s*(.*?)(?=Q:|$)'
        matches = re.findall(qa_pattern, text, re.DOTALL | re.IGNORECASE)

        if not matches:
            # If no Q&A pattern found, treat as single document
            logger.info(f"No Q&A pattern found in {filename}, treating as single document")
            return [Document(
                page_content=text.strip(),
                metadata={
                    'source': source,
                    'filename': filename,
                    'content_type': 'kb_page',
                    'is_faq': False
                }
            )]

        logger.info(f"Found {len(matches)} Q&A pairs in {filename}")

        for idx, (question, answer) in enumerate(matches, start=1):
            question = question.strip()
            answer = answer.strip()

            if question and answer:
                # Create combined content for better semantic search
                content = f"Q: {question}\n\nA: {answer}"

                doc = Document(
                    page_content=content,
                    metadata={
                        'source': source,
                        'filename': filename,
                        'content_type': 'kb_page',
                        'is_faq': True,
                        'question': question,
                        'answer': answer,
                        'qa_index': idx
                    }
                )
                documents.append(doc)

        return documents

    @staticmethod
    def parse_document(file_path: str) -> List[Dict[str, Any]]:
        """
        Parse any supported document format.
        Supports: TXT, MD, PDF, DOCX, PPTX, EPUB

        Args:
            file_path: Path to the document file

        Returns:
            List of document sections with metadata
        """
        ext = os.path.splitext(file_path)[1].lower()

        if ext in ['.txt', '.md']:
            return GenericKnowledgeBaseParser._parse_text(file_path)
        elif ext == '.pdf':
            return GenericKnowledgeBaseParser._parse_pdf(file_path)
        elif ext in ['.docx', '.doc']:
            return GenericKnowledgeBaseParser._parse_docx(file_path)
        elif ext in ['.pptx', '.ppt']:
            return GenericKnowledgeBaseParser._parse_pptx(file_path)
        elif ext == '.epub':
            return GenericKnowledgeBaseParser._parse_epub(file_path)
        else:
            logger.warning(f"Unsupported file format: {ext} for file {file_path}")
            return []

    @staticmethod
    def _parse_text(file_path: str) -> List[Dict[str, Any]]:
        """Parse plain text or markdown files."""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()

            filename = os.path.basename(file_path)

            # Detect if this looks like a FAQ document
            is_faq = bool(re.search(r'(Q:|Question:|FAQ|Frequently Asked)', content, re.IGNORECASE))

            return [{
                'content': content,
                'metadata': {
                    'source': file_path,
                    'filename': filename,
                    'file_type': 'text',
                    'content_type': 'kb_page',
                    'is_faq': is_faq
                }
            }]
        except Exception as e:
            logger.error(f"Error parsing text file {file_path}: {e}")
            return []

    @staticmethod
    def _parse_pdf(file_path: str) -> List[Dict[str, Any]]:
        """Parse PDF files."""
        try:
            from pypdf import PdfReader

            reader = PdfReader(file_path)
            filename = os.path.basename(file_path)
            sections = []

            for page_num, page in enumerate(reader.pages, start=1):
                text = page.extract_text()
                if text.strip():
                    sections.append({
                        'content': text,
                        'metadata': {
                            'source': file_path,
                            'filename': filename,
                            'file_type': 'pdf',
                            'content_type': 'kb_page',
                            'page_number': page_num
                        }
                    })

            logger.info(f"Parsed {len(sections)} pages from PDF: {filename}")
            return sections

        except ImportError:
            logger.error("pypdf not installed. Install with: pip install pypdf")
            return []
        except Exception as e:
            logger.error(f"Error parsing PDF {file_path}: {e}")
            return []

    @staticmethod
    def _parse_docx(file_path: str) -> List[Dict[str, Any]]:
        """Parse DOCX files."""
        try:
            from docx import Document as DocxDocument

            doc = DocxDocument(file_path)
            filename = os.path.basename(file_path)

            # Extract all paragraphs
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            content = '\n\n'.join(paragraphs)

            return [{
                'content': content,
                'metadata': {
                    'source': file_path,
                    'filename': filename,
                    'file_type': 'docx',
                    'content_type': 'kb_page'
                }
            }]

        except ImportError:
            logger.error("python-docx not installed. Install with: pip install python-docx")
            return []
        except Exception as e:
            logger.error(f"Error parsing DOCX {file_path}: {e}")
            return []

    @staticmethod
    def _parse_pptx(file_path: str) -> List[Dict[str, Any]]:
        """Parse PPTX files."""
        try:
            from pptx import Presentation

            prs = Presentation(file_path)
            filename = os.path.basename(file_path)
            sections = []

            for slide_num, slide in enumerate(prs.slides, start=1):
                text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_parts.append(shape.text)

                if text_parts:
                    content = '\n'.join(text_parts)
                    sections.append({
                        'content': content,
                        'metadata': {
                            'source': file_path,
                            'filename': filename,
                            'file_type': 'pptx',
                            'content_type': 'kb_page',
                            'slide_number': slide_num
                        }
                    })

            logger.info(f"Parsed {len(sections)} slides from PPTX: {filename}")
            return sections

        except ImportError:
            logger.error("python-pptx not installed. Install with: pip install python-pptx")
            return []
        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {e}")
            return []

    @staticmethod
    def _parse_epub(file_path: str) -> List[Dict[str, Any]]:
        """Parse EPUB files."""
        try:
            import ebooklib
            from ebooklib import epub
            from bs4 import BeautifulSoup

            book = epub.read_epub(file_path)
            filename = os.path.basename(file_path)
            sections = []

            for item in book.get_items():
                if item.get_type() == ebooklib.ITEM_DOCUMENT:
                    soup = BeautifulSoup(item.get_content(), 'html.parser')
                    text = soup.get_text()
                    if text.strip():
                        sections.append({
                            'content': text,
                            'metadata': {
                                'source': file_path,
                                'filename': filename,
                                'file_type': 'epub',
                                'content_type': 'kb_page'
                            }
                        })

            logger.info(f"Parsed {len(sections)} sections from EPUB: {filename}")
            return sections

        except ImportError:
            logger.error("ebooklib or beautifulsoup4 not installed. Install with: pip install ebooklib beautifulsoup4")
            return []
        except Exception as e:
            logger.error(f"Error parsing EPUB {file_path}: {e}")
            return []


class GenericKnowledgeBaseLoader:
    """Generic loader for knowledge base documents."""

    def __init__(self, directory: str):
        """
        Initialize the loader.

        Args:
            directory: Directory containing knowledge base files
        """
        self.directory = directory
        self.parser = GenericKnowledgeBaseParser()

    def load(self) -> List[Document]:
        """
        Load all documents from directory.

        Returns:
            List of Document objects
        """
        documents = []
        supported_extensions = ['.txt', '.md', '.pdf', '.docx', '.doc', '.pptx', '.ppt', '.epub']

        if not os.path.exists(self.directory):
            logger.warning(f"Directory does not exist: {self.directory}")
            return documents

        for root, _, files in os.walk(self.directory):
            for file in files:
                if any(file.endswith(ext) for ext in supported_extensions):
                    file_path = os.path.join(root, file)
                    try:
                        parsed = self.parser.parse_document(file_path)
                        for section in parsed:
                            doc = Document(
                                page_content=section['content'],
                                metadata=section['metadata']
                            )
                            documents.append(doc)
                    except Exception as e:
                        logger.error(f"Error loading document {file_path}: {e}")

        logger.info(f"Loaded {len(documents)} documents from {self.directory}")
        return documents
